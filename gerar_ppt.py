from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 1. Slide de Título
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Refatoração do Power BI para Web App"
    subtitle.text = "Validação de Dados e Comparativo de Telas\nApple Nordeste | Acrefort | 4E Equipamentos"
    
    empresas = [
        {
            "nome": "Apple Nordeste Comercio de Alimentos Ltda",
            "ebitda_pbi": "R$ 8.502.794,71", "ebitda_dash": "R$ 8.502.794,71",
            "lucro_pbi": "R$ 5.247.709,60", "lucro_dash": "R$ 5.247.709,60",
            "custo_pbi": "R$ 15.142.549,05", "custo_dash": "R$ 15.142.549,05"
        },
        {
            "nome": "ACREFORT INDUSTRIA E COMERCIO DE RAÇÕES LTDA",
            "ebitda_pbi": "R$ 8.149.808,13", "ebitda_dash": "R$ 8.149.808,13",
            "lucro_pbi": "R$ -1.525.184,87", "lucro_dash": "R$ -1.525.184,87",
            "custo_pbi": "R$ 26.243.588,71", "custo_dash": "R$ 26.243.588,71"
        },
        {
            "nome": "4E EQUIPAMENTOS PARA CAMINHOES LTDA",
            "ebitda_pbi": "R$ 572.944,39", "ebitda_dash": "R$ 572.944,39",
            "lucro_pbi": "R$ 719.731,06", "lucro_dash": "R$ 719.731,06",
            "custo_pbi": "R$ 7.512.673,69", "custo_dash": "R$ 7.512.673,69"
        }
    ]
    
    # Slides por empresa
    for emp in empresas:
        # Slide da Empresa (Validação de Dados)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = f"Validação: {emp['nome']}"
        
        # Add table
        rows = 4
        cols = 3
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        table.cell(0, 0).text = "Indicador"
        table.cell(0, 1).text = "Power BI (Semântico)"
        table.cell(0, 2).text = "Dashboard (BigQuery)"
        
        # Row 1
        table.cell(1, 0).text = "EBITDA"
        table.cell(1, 1).text = emp["ebitda_pbi"]
        table.cell(1, 2).text = emp["ebitda_dash"]
        
        # Row 2
        table.cell(2, 0).text = "Lucro Líquido"
        table.cell(2, 1).text = emp["lucro_pbi"]
        table.cell(2, 2).text = emp["lucro_dash"]
        
        # Row 3
        table.cell(3, 0).text = "Custo Total"
        table.cell(3, 1).text = emp["custo_pbi"]
        table.cell(3, 2).text = emp["custo_dash"]
    
    # Save presentation
    prs.save("Apresentacao_Final_v2.pptx")
    print("Apresentação gerada com sucesso!")

if __name__ == "__main__":
    create_presentation()
